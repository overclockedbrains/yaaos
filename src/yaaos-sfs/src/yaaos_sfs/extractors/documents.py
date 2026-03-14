"""Document extractors — DOCX, PPTX, XLSX, EPUB, RTF.

All dependencies are optional. If a library is missing, that extractor
simply won't be registered and those file types will be skipped.
"""

from __future__ import annotations

import logging
from pathlib import Path

log = logging.getLogger("yaaos-sfs")


def extract_docx(path: Path) -> str | None:
    """Extract text from DOCX using python-docx."""
    try:
        from docx import Document

        doc = Document(str(path))
        parts = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # Preserve heading structure
                if para.style and para.style.name and para.style.name.startswith("Heading"):
                    level = para.style.name.replace("Heading ", "").strip()
                    try:
                        prefix = "#" * int(level)
                        parts.append(f"{prefix} {text}")
                    except ValueError:
                        parts.append(text)
                else:
                    parts.append(text)

        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))

        return "\n\n".join(parts) if parts else None
    except Exception as e:
        log.debug(f"DOCX extraction failed for {path.name}: {e}")
        return None


def extract_pptx(path: Path) -> str | None:
    """Extract text from PPTX using python-pptx."""
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)

                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_texts.append(" | ".join(cells))

            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_texts.append(f"[Speaker notes: {notes}]")

            if slide_texts:
                parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))

        return "\n\n".join(parts) if parts else None
    except Exception as e:
        log.debug(f"PPTX extraction failed for {path.name}: {e}")
        return None


def extract_xlsx(path: Path) -> str | None:
    """Extract text from XLSX using openpyxl."""
    try:
        from openpyxl import load_workbook

        wb = load_workbook(str(path), read_only=True, data_only=True)
        parts = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_parts = [f"## Sheet: {sheet_name}"]
            row_count = 0

            for row in ws.iter_rows(values_only=True):
                if row_count >= 1000:  # Cap at 1000 rows per sheet
                    sheet_parts.append(f"... ({ws.max_row - 1000} more rows)")
                    break
                cells = [str(cell) if cell is not None else "" for cell in row]
                line = " | ".join(cells).strip()
                if line and line != "| " * len(cells):
                    sheet_parts.append(line)
                    row_count += 1

            if len(sheet_parts) > 1:  # More than just the header
                parts.append("\n".join(sheet_parts))

        wb.close()
        return "\n\n".join(parts) if parts else None
    except Exception as e:
        log.debug(f"XLSX extraction failed for {path.name}: {e}")
        return None


def extract_epub(path: Path) -> str | None:
    """Extract text from EPUB using ebooklib + html2text fallback."""
    try:
        import ebooklib
        from ebooklib import epub
        import re

        book = epub.read_epub(str(path), options={"ignore_ncx": True})
        parts = []

        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            content = item.get_content().decode("utf-8", errors="replace")
            # Strip HTML tags (simple approach, no extra deps)
            text = re.sub(r"<[^>]+>", " ", content)
            text = re.sub(r"\s+", " ", text).strip()
            if text and len(text) > 20:  # Skip near-empty chapters
                parts.append(text)

        return "\n\n".join(parts) if parts else None
    except Exception as e:
        log.debug(f"EPUB extraction failed for {path.name}: {e}")
        return None


def extract_rtf(path: Path) -> str | None:
    """Extract text from RTF using striprtf."""
    try:
        from striprtf.striprtf import rtf_to_text

        raw = path.read_text(encoding="utf-8", errors="replace")
        text = rtf_to_text(raw)
        return text.strip() if text and text.strip() else None
    except Exception as e:
        log.debug(f"RTF extraction failed for {path.name}: {e}")
        return None


def register_extractors() -> None:
    """Register document extractors for available libraries."""
    from . import register

    # Each extractor is registered only if its dependency is importable
    _optional_register(register, [".docx"], extract_docx, "python-docx", "docx")
    _optional_register(register, [".pptx"], extract_pptx, "python-pptx", "pptx")
    _optional_register(register, [".xlsx", ".xls"], extract_xlsx, "openpyxl", "openpyxl")
    _optional_register(register, [".epub"], extract_epub, "ebooklib", "ebooklib")
    _optional_register(register, [".rtf"], extract_rtf, "striprtf", "striprtf")


def _optional_register(register_fn, extensions, extractor, pkg_name, import_name):
    """Register an extractor only if its dependency can be imported."""
    try:
        __import__(import_name)
        register_fn(extensions, extractor)
        log.debug(f"Registered {pkg_name} extractor for {extensions}")
    except ImportError:
        log.debug(f"{pkg_name} not installed — skipping {extensions} support")
